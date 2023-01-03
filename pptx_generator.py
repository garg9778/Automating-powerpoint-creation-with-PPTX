import collections
import collections.abc
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pathlib import Path
import datetime
import pandas as pd
import mysql.connector as sql
import os
import warnings

warnings.filterwarnings('ignore')

def run_query(query, path=False, database='trends_mktplace', user='root', password=''):
    print('Retrieving the data from', query)
    if path:
        with open(query, 'r') as f:
            query = f.read()
            
    if password == '':
        if 'MYSQL_PASSWORD' in os.environ:
            password = os.environ['MYSQL_PASSWORD']
        else:
            raise ValueError('Password must be given or invalid')
    
    with sql.connect(host='localhost',
                     database=database,
                     user=user,
                     password=password) as c:
        
        return pd.read_sql(query, c)

    
def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = str(new_text)


def edit_group_shapes_on_page(page: dict, group_keyword: str, sub_group_keys_values: dict,
                              text_filed_pattern: dict):
    found = False
    for shape in page:
        if group_keyword in shape and shape[-1].isnumeric():
            found = True
            for sub_group_key in sub_group_keys_values:
                assert sub_group_key in page[shape], f'"{sub_group_key}" shape not found for {shape}'
            
            group_shape = page[shape]
            group_number = int(shape[-1])
            
            for sub_group_key in sub_group_keys_values:
                para = group_shape[sub_group_key].text_frame.paragraphs[0]
                val_list = sub_group_keys_values[sub_group_key]
                if sub_group_key in text_filed_pattern:
                    pattern = text_filed_pattern[sub_group_key]
                    new_val = pattern.format(val_list[group_number - 1])
                else:
                    new_val = val_list[group_number - 1]
                replace_paragraph_text_retaining_initial_formatting(para, new_val)
            
    if not found:
        raise AssertionError(f'"{group_keyword}" not found on the given page')


# Path for all

parent = Path(__file__).parent
sql_folder_path = parent / 'sql_scripts'
output_folder_path = parent / 'generated_pptx'
template_folder_path = parent / 'template'

template_file_path = template_folder_path / 'template.pptx'

qr_s1_line_chart = sql_folder_path / 'query_1_line_chart.sql'
qr_s1_product_name = sql_folder_path / 'query_1_products.sql'
qr_s2_table = sql_folder_path / 'query_2.sql'
qr_s3_line_chart = sql_folder_path / 'query_3_line_chart.sql'
qr_s3_pie_chart = sql_folder_path / 'query_3_pie_chart.sql'
qr_s4_bar_chart = sql_folder_path / 'query_4.sql'
qr_s4_region = sql_folder_path / 'query_4_regions.sql'


if __name__ == '__main__':
    print('Getting pptx template from ', template_file_path)
    prs = Presentation(template_file_path)

    slides = prs.slides

    # Create page dict
    pages = {}
    for i, slide in enumerate(slides):
        shape_dict = {}
        for shape in slide.shapes:      
            
            if shape.shape_type == 6:  # group_shape
                group_dict = {}
                for g_shape in shape.shapes:
                    group_dict[g_shape.name] = g_shape
                shape_dict[shape.name] = group_dict
                
            elif shape.shape_type == 3:  # chart
                shape_dict[shape.name] = shape.chart

            else:
                shape_dict[shape.name] = shape
            
            pages[i+1] = shape_dict


    ### Page 2

    # Get the data
    monthly_sales = run_query(qr_s1_line_chart, path=True).fillna(0)
    product_names = run_query(qr_s1_product_name, path=True)['product_line']

    sales_current_month_top_3 = monthly_sales.iloc[-1,1:].values
    
    # Modify line chart
    # Note that the line chart must be named as 'line_chart' on page 2
    line_chart_data = ChartData()
    line_chart_data.categories = monthly_sales.iloc[:,0].values


    for col in range(1, len(monthly_sales.columns)):
        line_chart_data.add_series(monthly_sales.columns[col], monthly_sales.iloc[:,col].values)

    line_chart = pages[2]['line_chart']
    line_chart.replace_data(line_chart_data)
        

    # Change product name and number
    # Note that each group must be named as 'product_X' 
    # and the shapes inside each group shape must be name 'name' and 'sales' respectively

    new_data_dict = {'name': product_names,
                     'sales': sales_current_month_top_3}

    pattern_dict = {'sales': "{:,.0f}"}
    
    edit_group_shapes_on_page(page=pages[2], group_keyword='product',
                              sub_group_keys_values=new_data_dict,
                              text_filed_pattern= pattern_dict)


    # Change the monthly total sales change
    prev_month_top_3 = monthly_sales.iloc[-2,1:].sum()
    cur_month_top_3 = monthly_sales.iloc[-1,1:].sum()

    percent_change = ((cur_month_top_3 - prev_month_top_3) / prev_month_top_3 * 100).round().astype(int)
    sign = '+' if percent_change >= 0 else '' 

    summary = pages[2]['summary']
    new_percent_text = f'{sign} {percent_change}%'
    para_name = summary['percent'].text_frame.paragraphs[0]
    summary['percent'] = replace_paragraph_text_retaining_initial_formatting(para_name, new_percent_text)



    ### Page 3
    df3 = run_query(qr_s2_table, path=True)

    table_3 = pages[3]['table'].table

    for row in range(df3.shape[0]):
        for col in range(df3.shape[1]):
            cell = table_3.cell(row + 1, col)
            data = df3.iloc[row, col]
            if isinstance(data, str):
                f_data = data
            else:
                f_data = "${:,}".format(data)
                
            cell.text_frame.text = f_data
            cell_font = cell.text_frame.paragraphs[0].font  # get font object

            cell_font.size = Pt(15) if col != 0 else Pt(21)
            cell_font.bold = True
            cell_font.color.rgb = RGBColor(0, 0, 0) if col != 0 else RGBColor(255, 255, 255)
            cell_font.name = 'Fira Sans Extra Condensed'
            
    ### Page 4

    # Get the data
    cur_month_sales = run_query(qr_s3_pie_chart, path=True)
    product_line = cur_month_sales['product_line'].values
    percentages = (cur_month_sales['sale'] * 100 / cur_month_sales['sale'].sum()).round().astype(int).values
    annual_sales_6_products = run_query(qr_s3_line_chart, path=True).fillna(0)

    # Modify pie chart
    # Note that the pie chart must be named as 'pie_chart' on page 4
    pie_chart_data = ChartData()
    pie_chart_data.categories = product_line
    pie_chart_data.add_series('products', percentages)

    pie_chart = pages[4]['pie_chart']
    pie_chart.replace_data(pie_chart_data)

    # Update product sales percentage
    new_data_dict = {'name': product_line,
                     'percent': percentages}

    pattern_dict = {'sales': "{:,.0f}%"}
    
    edit_group_shapes_on_page(page=pages[4], group_keyword='product',
                              sub_group_keys_values=new_data_dict,
                              text_filed_pattern= pattern_dict)
            
    # Modify line chart
    # Note that the line chart must be named as 'line_chart' on page 2
    line_chart_data = ChartData()
    line_chart_data.categories = annual_sales_6_products.iloc[:,0].values

    chart_df = annual_sales_6_products
    for col in range(1, chart_df.shape[1]):
        line_chart_data.add_series(chart_df.columns[col], chart_df.iloc[:,col].values)

    line_chart = pages[4]['line_chart']
    line_chart.replace_data(line_chart_data)

    #### P5
    annual_sales_region = run_query(qr_s4_bar_chart, path=True).fillna(0)
    region_names = run_query(qr_s4_region, path=True)['region']

    # Modify bar chart
    # Note that the line chart must be named as 'line_chart' on page 2
    bar_chart_data = ChartData()
    bar_chart_data.categories = annual_sales_region.iloc[:,0].values


    for col in range(1, len(annual_sales_region.columns)):
        bar_chart_data.add_series(region_names[col - 1], annual_sales_region.iloc[:,col].values)
        
    bar_chart = pages[5]['bar_chart']
    bar_chart.replace_data(bar_chart_data)

    # Modify max date
    date_r1, max_r1, _ = annual_sales_region[annual_sales_region['R1'] == annual_sales_region['R1'].max()].values[0]
    date_r2, _, max_r2 = annual_sales_region[annual_sales_region['R2'] == annual_sales_region['R2'].max()].values[0]

    new_data_dict = {'date': [date_r1, date_r2],
                     'sales': [max_r1, max_r2]}

    pattern_dict = {'sales': "{:,.0f}"}
    
    edit_group_shapes_on_page(page=pages[5], group_keyword='best_month',
                              sub_group_keys_values=new_data_dict,
                              text_filed_pattern= pattern_dict)

    # Modify region
    for i in range(1, 3):
        region = pages[5][f'region_{i}']
        para = region.text_frame.paragraphs[0]
        new_region_name = region_names[i - 1]
        replace_paragraph_text_retaining_initial_formatting(para, new_region_name)

    # Save pptx
    time = datetime.datetime.now().strftime("%m%d%y_%H%M")
    file_name = 'generated_pptx'
    output_filename = f'{file_name}_{time}.pptx'
    output_file_path = output_folder_path / output_filename
    print('Creating a new pptx ...')
    prs.save(output_file_path)
    print('New file has been created as', output_file_path)

    # Ask to open file
    open_needed = input('Do you want to open the generated pptx? (Y/N): ')
    if open_needed.lower() in ['y', 'yes']:
        os.startfile(output_file_path)
